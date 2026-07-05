"""Behavioral smoke test for v1.6.1: the write_pptx 模板视觉升级.

v1.6.1 redesigns the PowerPoint template's *look* (the把关人 opens the export in PowerPoint COM and
inspects it by eye). This test builds ONE real sample deck exercising every layout the redesign
touches and asserts the load-bearing structural facts back through python-pptx:

  样例 deck (6 slides):
    1. title   — full-bleed dark COVER + subtitle + date.
    2. content — 3 bullets (large, vertically-centred single column).
    3. content — 8 bullets (auto two-column: TWO body textboxes).
    4. stats   — 4 number cards.
    5. table   — the classic styled table (back-compat).
    6. closing — full-bleed dark closing.

  Assertions:
    * cover slide has a full-bleed rectangle filled with the token `title_bg` (== primary),
    * cover carries NO page-number footer, closing carries NO page-number footer,
    * the 8-bullet content slide produced exactly TWO body textboxes (two-column split),
    * the stats slide has exactly 4 card panels (fill == token `card_fill`),
    * a `closing` slide exists and shows '谢谢',
    * the deck reads back as a valid 6-slide pptx,
    * back-compat: empty items / >6 items / >10 bullets all behave (error / truncate).

The sample .pptx absolute path is printed at the end so the把关人 can COM-export it for the eye check.

Run with UTF-8:  python -X utf8 tests/smoke_v161.py
"""

import os
import sys
import tempfile
import zipfile

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_SRC = os.path.join(_ROOT, "src")
sys.path.insert(0, _SRC)

_DATA = os.path.join(tempfile.gettempdir(), "acc_smoke_v161_data")
os.makedirs(_DATA, exist_ok=True)
os.environ["WCW_DATA_DIR"] = _DATA

import ai_computer_control.server as server  # noqa: E402
from ai_computer_control.tools import office_style as tokens  # noqa: E402

_FNS = {t.name: t.fn for t in server.mcp._tool_manager.list_tools()}
_FAILURES: list[str] = []


def check(cond: bool, msg: str):
    print(f"  [{'ok  ' if cond else 'FAIL'}] {msg}")
    if not cond:
        _FAILURES.append(msg)


def _shape_fill_rgb(sh):
    """Return the shape's solid fill as an (r,g,b) tuple, or None if it has no readable solid fill."""
    try:
        rgb = sh.fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def main() -> int:
    from pptx import Presentation
    from pptx.util import Emu

    biz = tokens.get_style("business")
    title_bg = tokens.rgb_tuple(biz["title_bg"])
    card_fill = tokens.rgb_tuple(biz["card_fill"])

    print("== registry: write_pptx present ==")
    check("write_pptx" in _FNS, "write_pptx present in the tool registry")

    # New v1.6.1 tokens exist on every style.
    print("\n== tokens: title_bg / on_dark_subtle / card_fill on all three styles ==")
    for st in ("business", "minimal", "vibrant"):
        t = tokens.get_style(st)
        for k in ("title_bg", "on_dark_subtle", "card_fill"):
            check(k in t and isinstance(t[k], str) and len(t[k]) == 6,
                  f"{st}.{k} present as 6-hex (got {t.get(k)!r})")

    # ---------------------------------------------------------------- build the sample deck
    print("\n== build the v1.6.1 sample deck (cover + 3-bullet + 8-bullet + stats + table + closing) ==")
    pptx_path = os.path.join(_DATA, "如意_v161_视觉升级样例.pptx")
    if os.path.exists(pptx_path):
        os.remove(pptx_path)
    slides = [
        {"type": "title", "title": "季度业绩汇报",
         "subtitle": "2026 财年第一季度 · 商务简约", "date": "2026-07-05"},
        {"type": "content", "title": "核心要点（三条·大字居中）", "bullets": [
            "销售额同比增长 26%，创同期新高",
            {"text": "华东区连续三季领跑", "level": 1},
            "第四季度冲刺目标已锁定"]},
        {"type": "content", "title": "推进事项（八条·自动两栏）", "bullets": [
            "完成年度预算复盘", "上线新版客户门户", "华南区团队扩编",
            "供应链二级备份", "季度 OKR 对齐会", "老客户续约专项",
            "数据看板 2.0", "合规审计整改"]},
        {"type": "stats", "title": "关键指标", "items": [
            {"label": "营业收入", "value": "¥18.9M", "note": "同比 +33%"},
            {"label": "新增客户", "value": "1,240", "note": "环比 +12%"},
            {"label": "毛利率", "value": "41.2%"},
            {"label": "客户 NPS", "value": "62", "note": "行业前 10%"}]},
        {"type": "table", "title": "季度明细", "headers": ["季度", "销售额", "增长"],
         "rows": [["第一季度", "1250", "12%"], ["第二季度", "1580", "26%"],
                  ["第三季度", "1420", "-10%"], ["第四季度", "1890", "33%"]]},
        {"type": "closing", "title": "谢谢", "subtitle": "欢迎交流与指正"},
    ]
    wp = _FNS["write_pptx"](path=pptx_path, slides=slides, style="business")
    check(isinstance(wp, dict) and wp.get("success") is True, f"write_pptx ok (got {wp})")
    check(wp.get("slides") == 6, f"reported 6 slides (got {wp.get('slides')})")
    check(wp.get("style") == "business", f"reported style business (got {wp.get('style')})")
    check(os.path.exists(pptx_path), "sample pptx exists on disk")
    check(zipfile.is_zipfile(pptx_path), "sample pptx is a valid OOXML zip container")

    prs = Presentation(pptx_path)
    check(len(prs.slides) == 6, f"python-pptx reads back 6 slides (got {len(prs.slides)})")

    slide_w, slide_h = prs.slide_width, prs.slide_height
    # v1.7: content body area top was raised 1.75in -> 1.4in (把关人: content 重心略低). Lower this
    # bound to match so the body-textbox detection still brackets the (now higher) content boxes.
    body_lo = Emu(int(1.35 * 914400))  # just below the raised content body top (title bar is 1.25in tall)
    body_hi = Emu(int(6.95 * 914400))  # above the footer

    # --- cover: full-bleed rectangle filled with title_bg ---
    cover = prs.slides[0]
    cover_bg = any(
        _shape_fill_rgb(sh) == title_bg and sh.width == slide_w and sh.height == slide_h
        for sh in cover.shapes)
    check(cover_bg, f"cover has a full-bleed background filled with title_bg {biz['title_bg']}")
    # cover subtitle + date text present
    cover_text = "\n".join(sh.text_frame.text for sh in cover.shapes if sh.has_text_frame)
    check("季度业绩汇报" in cover_text, "cover title text present")
    check("2026-07-05" in cover_text, "cover date rendered")

    # --- cover carries no page-number footer -----------------------------------------------
    # The footer is the only bottom-right textbox whose text is just the page number '1'.
    def _has_page_footer(slide, page_no):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == str(page_no) \
                    and sh.top and sh.top > Emu(int(6.9 * 914400)):
                return True
        return False
    check(not _has_page_footer(cover, 1), "cover has NO page-number footer")

    # --- 8-bullet content -> two body textboxes (two-column) -------------------------------
    twocol = prs.slides[2]
    body_boxes = [sh for sh in twocol.shapes
                  if sh.has_text_frame and sh.top and body_lo < sh.top < body_hi
                  and sh.text_frame.text.strip()]
    check(len(body_boxes) == 2, f"8-bullet content split into TWO body textboxes (got {len(body_boxes)})")

    # --- 3-bullet content -> single body textbox -------------------------------------------
    onecol = prs.slides[1]
    body_boxes1 = [sh for sh in onecol.shapes
                   if sh.has_text_frame and sh.top and body_lo < sh.top < body_hi
                   and sh.text_frame.text.strip()]
    check(len(body_boxes1) == 1, f"3-bullet content is a SINGLE body textbox (got {len(body_boxes1)})")

    # --- stats: exactly 4 card panels (fill == card_fill) ----------------------------------
    stats = prs.slides[3]
    card_panels = [sh for sh in stats.shapes if _shape_fill_rgb(sh) == card_fill]
    check(len(card_panels) == 4, f"stats slide has exactly 4 card panels (got {len(card_panels)})")
    stats_text = "\n".join(sh.text_frame.text for sh in stats.shapes if sh.has_text_frame)
    check("¥18.9M" in stats_text and "同比 +33%" in stats_text, "stats value + note text present")

    # --- closing: exists, full-bleed dark, shows 谢谢, no footer ----------------------------
    closing = prs.slides[5]
    closing_bg = any(
        _shape_fill_rgb(sh) == title_bg and sh.width == slide_w and sh.height == slide_h
        for sh in closing.shapes)
    check(closing_bg, "closing slide has a full-bleed title_bg background")
    closing_text = "\n".join(sh.text_frame.text for sh in closing.shapes if sh.has_text_frame)
    check("谢谢" in closing_text, "closing slide shows 谢谢")
    check(not _has_page_footer(closing, 6), "closing has NO page-number footer")

    # --- content slide (non-cover) DOES carry a footer (back-compat) -----------------------
    check(_has_page_footer(onecol, 2), "content slide keeps its page-number footer (back-compat)")

    # ---------------------------------------------------------------- error / edge paths
    print("\n== error & edge paths: empty stats / >6 stats / >10 bullets ==")
    e_empty = _FNS["write_pptx"](path=os.path.join(_DATA, "e_empty.pptx"),
                                 slides=[{"type": "stats", "title": "x", "items": []}])
    check(isinstance(e_empty, dict) and e_empty.get("ok") is False and bool(e_empty.get("error")),
          f"empty stats items -> 人话 error (got {e_empty})")
    e_many = _FNS["write_pptx"](path=os.path.join(_DATA, "e_many.pptx"),
                                slides=[{"type": "stats", "title": "x",
                                         "items": [{"label": "a", "value": "1"}] * 7}])
    check(isinstance(e_many, dict) and e_many.get("ok") is False and bool(e_many.get("error")),
          f">6 stats items -> 人话 error (got {e_many})")
    trunc_path = os.path.join(_DATA, "trunc.pptx")
    tr = _FNS["write_pptx"](path=trunc_path,
                            slides=[{"type": "content", "title": "太多条",
                                     "bullets": ["条目" + str(i) for i in range(15)]}])
    check(isinstance(tr, dict) and tr.get("success") is True, "15-bullet content still generates (no crash)")
    tprs = Presentation(trunc_path)
    ttext = "\n".join(sh.text_frame.text for sl in tprs.slides for sh in sl.shapes if sh.has_text_frame)
    check("内容过多" in ttext, ">10 bullets truncated with 「…（内容过多，建议拆页）」 note")

    # ---------------------------------------------------------------- three-style coverage
    print("\n== all three styles render the cover + stats + closing without error ==")
    tri_slides = [
        {"type": "title", "title": "封面", "subtitle": "副标题", "date": "2026-07-05"},
        {"type": "stats", "title": "指标", "items": [
            {"label": f"K{i}", "value": str(i)} for i in range(1, 7)]},
        {"type": "closing"},
    ]
    for st in ("minimal", "vibrant", "business"):
        r = _FNS["write_pptx"](path=os.path.join(_DATA, f"tri_{st}.pptx"),
                               slides=tri_slides, style=st)
        check(isinstance(r, dict) and r.get("success") is True and r.get("style") == st,
              f"style '{st}' renders cover+stats(6)+closing ok (got {r.get('error') or r.get('style')})")

    print()
    if _FAILURES:
        print(f"FAILED: {len(_FAILURES)} assertion(s)")
        for f in _FAILURES:
            print("  -", f)
        print("ACC-V161 SMOKE: FAIL")
        return 1
    print("样例 pptx（把关人用 PowerPoint COM 导图亲验）:")
    print("  pptx :", os.path.abspath(pptx_path))
    print("ACC-V161 SMOKE: ALL PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
