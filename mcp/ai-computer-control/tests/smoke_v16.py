"""Behavioral smoke test for v1.6: the Office 模板驱动 tools.

Exercises all four new tools end-to-end, producing REAL files and reading key style points back:

  (T1) excel_beautify — write_excel a Chinese dataset, beautify it, reopen and assert:
        header fill == the design-token primary colour, header font bold + white, freeze_panes set,
        auto_filter set, a data row carries the zebra tint. Run it TWICE and assert the second pass
        produces the same styling (idempotent — no accumulation).
  (T2) excel_chart — insert a bar chart; reopen and assert the sheet has >=1 chart. Also assert a
        malformed data_range degrades to an {error} (defensive path).
  (T3) write_pptx — one slide of each of the 4 types (title/content/table/image) with Chinese text;
        assert file exists, is a valid zip/pptx, slide count == 4, and the title text is present.
        image slide reuses the .png from T4.
  (T4) chart_image — render a Chinese-titled bar chart .png; assert file exists, > 10 KB, PNG magic,
        and font is a real CJK font (not the no-font warning path — this box ships msyh.ttc).
        Also assert a multi-series pie degrades to an {error}.
  (deg) python-pptx / matplotlib MISSING -> the tools degrade to an install-guiding {error} without
        crashing (simulated in child processes that block the import).

All output goes to a throwaway data dir. Prints per-check PASS/FAIL and exits non-zero on any failure.

Run with UTF-8:  python -X utf8 tests/smoke_v16.py
"""

import os
import subprocess
import sys
import tempfile
import zipfile

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_SRC = os.path.join(_ROOT, "src")
sys.path.insert(0, _SRC)

_DATA = os.path.join(tempfile.gettempdir(), "acc_smoke_v16_data")
os.makedirs(_DATA, exist_ok=True)
os.environ["WCW_DATA_DIR"] = _DATA

import ai_computer_control.server as server  # noqa: E402
from ai_computer_control.tools import office_style as tokens  # noqa: E402

_FNS = {t.name: t.fn for t in server.mcp._tool_manager.list_tools()}
_FAILURES: list[str] = []


def check(cond: bool, msg: str):
    print(f"  [{'ok  ' if cond else 'FAIL'}] {msg}")
    if not cond:
        _FAILURES.append(msg)


def _norm(c: str) -> str:
    """openpyxl stores colours as 8-hex ARGB uppercase; normalize a token hex the same way."""
    return tokens.argb(c).upper()


def main() -> int:
    biz = tokens.get_style("business")

    print("== registry: all four v1.6 tools present ==")
    for name in ("excel_beautify", "excel_chart", "write_pptx", "chart_image"):
        check(name in _FNS, f"{name} present in the tool registry")

    # ---------------------------------------------------------------- T1 excel_beautify
    print("\n== (T1) excel_beautify: write a Chinese dataset, beautify, read style back ==")
    xlsx = os.path.join(_DATA, "季度业绩汇报.xlsx")
    if os.path.exists(xlsx):
        os.remove(xlsx)
    headers = ["季度", "销售额(万元)", "同比增长", "负责人"]
    data = [
        ["第一季度", 1250, "12%", "张伟"],
        ["第二季度", 1580, "26%", "李娜"],
        ["第三季度", 1420, "-10%", "王强"],
        ["第四季度", 1890, "33%", "赵敏"],
    ]
    w = _FNS["write_excel"](path=xlsx, data=data, headers=headers, sheet_name="业绩")
    check(isinstance(w, dict) and w.get("success") is True, f"write_excel wrote the base file (got {w.get('error')})")

    b = _FNS["excel_beautify"](path=xlsx, style="business", header_row=1)
    check(isinstance(b, dict) and b.get("success") is True, f"excel_beautify ok (got {b})")
    check(b.get("style") == "business", f"reported style == business (got {b.get('style')})")

    from openpyxl import load_workbook
    wb = load_workbook(xlsx)
    ws = wb["业绩"]
    hc = ws.cell(row=1, column=1)
    check(hc.fill.fgColor.rgb == _norm(biz["header_fill"]),
          f"header fill == token primary {_norm(biz['header_fill'])} (got {hc.fill.fgColor.rgb})")
    check(bool(hc.font.bold), "header font is bold")
    check(str(hc.font.color.rgb).upper() == _norm(biz["header_font_color"]),
          f"header font colour is white (got {hc.font.color.rgb})")
    check(ws.freeze_panes == "A2", f"freeze_panes frozen below header (got {ws.freeze_panes})")
    check(bool(ws.auto_filter.ref), f"auto_filter set (got {ws.auto_filter.ref})")
    # zebra: 2nd data row (row 3) should carry the zebra tint; 1st data row (row 2) should not.
    r3 = ws.cell(row=3, column=1).fill.fgColor.rgb
    check(r3 == _norm(biz["zebra_fill"]), f"2nd data row has zebra tint {_norm(biz['zebra_fill'])} (got {r3})")
    # numeric column (col 2) right-aligned
    check(ws.cell(row=2, column=2).alignment.horizontal == "right", "numeric column right-aligned")
    wb.close()

    # idempotency: run again, styling must be identical (not doubled/shifted).
    b2 = _FNS["excel_beautify"](path=xlsx, style="business", header_row=1)
    check(isinstance(b2, dict) and b2.get("success") is True, "excel_beautify second pass ok (idempotent)")
    wb2 = load_workbook(xlsx)
    ws2 = wb2["业绩"]
    check(ws2.cell(row=1, column=1).fill.fgColor.rgb == _norm(biz["header_fill"]),
          "idempotent: header fill unchanged after 2nd pass")
    check(ws2.cell(row=3, column=1).fill.fgColor.rgb == _norm(biz["zebra_fill"]),
          "idempotent: zebra tint unchanged after 2nd pass")
    check(ws2.freeze_panes == "A2", "idempotent: freeze_panes unchanged after 2nd pass")
    wb2.close()

    # ---------------------------------------------------------------- T2 excel_chart
    print("\n== (T2) excel_chart: insert a bar chart; malformed range degrades ==")
    ch = _FNS["excel_chart"](path=xlsx, sheet="业绩", chart_type="bar",
                             data_range="A1:B5", title="季度销售额", target_cell="H2")
    check(isinstance(ch, dict) and ch.get("success") is True, f"excel_chart bar inserted (got {ch})")
    wb3 = load_workbook(xlsx)
    ws3 = wb3["业绩"]
    n_charts = len(ws3._charts) if hasattr(ws3, "_charts") else 0
    check(n_charts >= 1, f"sheet has >=1 chart after insert (got {n_charts})")
    wb3.close()
    bad = _FNS["excel_chart"](path=xlsx, sheet="业绩", chart_type="bar",
                             data_range="NONSENSE", title="x")
    check(isinstance(bad, dict) and bad.get("ok") is False and bool(bad.get("error")),
          f"malformed data_range -> error (got {bad})")
    bad_t = _FNS["excel_chart"](path=xlsx, sheet="业绩", chart_type="donut",
                               data_range="A1:B5", title="x")
    check(isinstance(bad_t, dict) and bad_t.get("ok") is False,
          f"unknown chart_type -> error (got {bad_t})")

    # ---------------------------------------------------------------- T4 chart_image (before pptx, reused)
    print("\n== (T4) chart_image: Chinese bar chart .png, > 10 KB, CJK font, pie validation ==")
    png = os.path.join(_DATA, "季度销售图.png")
    if os.path.exists(png):
        os.remove(png)
    ci = _FNS["chart_image"](
        path=png, chart_type="bar", title="季度业绩汇报（万元）",
        data={"labels": ["第一季度", "第二季度", "第三季度", "第四季度"],
              "series": [{"name": "销售额", "values": [1250, 1580, 1420, 1890]},
                         {"name": "目标", "values": [1200, 1500, 1500, 1800]}]},
        style="business")
    check(isinstance(ci, dict) and ci.get("success") is True, f"chart_image ok (got {ci})")
    check(os.path.exists(png), "chart .png exists on disk")
    size = os.path.getsize(png) if os.path.exists(png) else 0
    check(size > 10 * 1024, f"chart .png > 10 KB (got {size} bytes)")
    magic_ok = False
    if os.path.exists(png):
        with open(png, "rb") as fh:
            magic_ok = fh.read(8) == b"\x89PNG\r\n\x1a\n"
    check(magic_ok, "chart .png has the PNG magic header")
    check(ci.get("font") and "warning" not in ci,
          f"a real CJK font resolved, no font warning (font={ci.get('font')}, warning={ci.get('warning')})")
    # pie single-series validation
    pie_bad = _FNS["chart_image"](path=os.path.join(_DATA, "pie.png"), chart_type="pie",
                                 title="x", data={"labels": ["a", "b"],
                                 "series": [{"name": "s1", "values": [1, 2]},
                                            {"name": "s2", "values": [3, 4]}]})
    check(isinstance(pie_bad, dict) and pie_bad.get("ok") is False,
          f"multi-series pie -> error (got {pie_bad})")

    # ---------------------------------------------------------------- T3 write_pptx
    print("\n== (T3) write_pptx: 4 slide types with Chinese, valid pptx, slide count == 4 ==")
    pptx_path = os.path.join(_DATA, "季度汇报.pptx")
    if os.path.exists(pptx_path):
        os.remove(pptx_path)
    slides = [
        {"type": "title", "title": "季度业绩汇报", "subtitle": "2026 财年第一季度 · 商务简约"},
        {"type": "content", "title": "核心要点", "bullets": [
            "销售额同比增长 26%", {"text": "华东区表现突出", "level": 1},
            "第四季度创历史新高"]},
        {"type": "table", "title": "季度明细", "headers": ["季度", "销售额", "增长"],
         "rows": [["第一季度", "1250", "12%"], ["第二季度", "1580", "26%"],
                  ["第三季度", "1420", "-10%"]]},
        {"type": "image", "title": "销售趋势图", "image_path": png, "caption": "图 1：季度销售额对比"},
    ]
    wp = _FNS["write_pptx"](path=pptx_path, slides=slides, style="business")
    check(isinstance(wp, dict) and wp.get("success") is True, f"write_pptx ok (got {wp})")
    check(wp.get("slides") == 4, f"reported 4 slides (got {wp.get('slides')})")
    check(os.path.exists(pptx_path), "pptx exists on disk")
    check(zipfile.is_zipfile(pptx_path), "pptx is a valid OOXML zip container")
    # Reopen with python-pptx: 4 slides, Chinese title text present.
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        check(len(prs.slides) == 4, f"python-pptx reads back 4 slides (got {len(prs.slides)})")
        all_text = []
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.has_text_frame:
                    all_text.append(sh.text_frame.text)
        joined = "\n".join(all_text)
        check("季度业绩汇报" in joined, "title slide Chinese text present on readback")
        check("核心要点" in joined and "季度明细" in joined, "content + table slide titles present")
    except Exception as e:  # noqa: BLE001
        check(False, f"pptx readback raised: {e}")
    # empty slides -> error
    empty = _FNS["write_pptx"](path=os.path.join(_DATA, "empty.pptx"), slides=[])
    check(isinstance(empty, dict) and empty.get("ok") is False, f"empty slides -> error (got {empty})")

    # ---------------------------------------------------------------- degrade paths
    print("\n== (deg) python-pptx / matplotlib MISSING -> install-guiding error, no crash ==")
    for mod, tool, kwargs in (
        ("pptx", "write_pptx", "path=os.path.join(d,'x.pptx'), slides=[{'type':'title','title':'x'}]"),
        ("matplotlib", "chart_image",
         "path=os.path.join(d,'x.png'), chart_type='bar', title='x', "
         "data={'labels':['a'],'series':[{'name':'s','values':[1]}]}"),
    ):
        child = (
            "import os,sys,tempfile\n"
            "class B:\n"
            f"  def find_spec(self,n,p=None,t=None):\n"
            f"    if n=={mod!r} or n.startswith({mod!r}+'.'): raise ImportError('blocked')\n"
            "    return None\n"
            "sys.meta_path.insert(0,B())\n"
            f"sys.modules.pop({mod!r},None)\n"
            "os.environ['WCW_DATA_DIR']=tempfile.mkdtemp(prefix='acc_v16_no_')\n"
            "d=os.environ['WCW_DATA_DIR']\n"
            "import ai_computer_control.server as s\n"
            "fns={t.name:t.fn for t in s.mcp._tool_manager.list_tools()}\n"
            f"r=fns[{tool!r}]({kwargs})\n"
            "e=r.get('error','')\n"
            f"ok=(r.get('ok') is False and bool(e) and ({mod!r} in e) and ('installer' in e or 'pip install' in e))\n"
            "print('DEGRADE_OK' if ok else 'DEGRADE_FAIL'); print('ERR='+str(e))\n"
        )
        env = dict(os.environ)
        env["PYTHONPATH"] = _SRC + os.pathsep + env.get("PYTHONPATH", "")
        env["PYTHONIOENCODING"] = "utf-8"
        proc = subprocess.run([sys.executable, "-X", "utf8", "-c", child],
                              capture_output=True, text=True, encoding="utf-8", timeout=180, env=env)
        out = proc.stdout or ""
        ok = "DEGRADE_OK" in out
        check(ok, f"{tool} without {mod} -> ok:False + install guidance (no crash)")
        if not ok:
            print("     child stdout:", out.strip()[:400])
            print("     child stderr:", (proc.stderr or "").strip()[:300])

    print()
    if _FAILURES:
        print(f"FAILED: {len(_FAILURES)} assertion(s)")
        for f in _FAILURES:
            print("  -", f)
        print("ACC-V16 SMOKE: FAIL")
        return 1
    print("产物路径:")
    print("  xlsx :", xlsx)
    print("  pptx :", pptx_path)
    print("  png  :", png)
    print("ACC-V16 SMOKE: ALL PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
