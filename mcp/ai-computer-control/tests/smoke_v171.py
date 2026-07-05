"""Behavioral smoke test for v1.7.1 — 图表轴标题 + minimal/vibrant 模板版式个性.

用户反馈两条:
  1. 「Excel 的图表好多都没有 X/Y 轴单位」 — excel_chart / chart_image 无轴标题。
  2. 「方案二(墨白极简)/方案三(活力现代)的模板可以再好看一点」 — 2/3 只是"换了配色的方案一"。

本冒烟断言 v1.7.1 的结构事实:

  ① excel_chart 轴标题:
     * 自动推导: x_title ← data_range 首列表头; y_title ← 单系列的系列表头 / 多系列留空。
     * 显式传参覆盖; '' 显式关闭; 饼图忽略。
     * delTitle 陷阱: x_axis.delete == False 且 y_axis.delete == False (否则轴连标题一起隐藏)。
     * 轴标题富文本字体 == 令牌 body_font (微软雅黑)。
     * 向后兼容: 老签名 (不带 x_title/y_title) 照常成功。
  ② chart_image 轴标签: 顶层参数 / data 内键 / 自动推导 (单系列 name) / '' 关闭 / 饼图忽略。
  ③ minimal PPT 版式: 封面底=白 (不再深底满铺) + 左侧青色竖线存在; content 无主色横幅;
     stats 无填充卡片 (无 card_fill 面板)。
  ④ vibrant PPT 版式: stats 卡片 prstGeom == roundRect (真圆角, adj≈0.12); 封面深靛底 + 几何装饰
     (圆/圆角矩形数量>=3); content 标题栏右端珊瑚圆点。
  ⑤ Word 版式: minimal 封面白底 (无深底题块) + 标题左侧竖线 (w:pBdr/w:left); vibrant 封面珊瑚粗条
     (0.3in exact row) + 深靛题块仍在。
  ⑥ Excel 表头版式: minimal 表头无满底色 + 墨黑粗体 + 底部 medium 青色下边框 (write_excel 落笔与
     excel_beautify 双通道一致); business/vibrant 表头仍主色满底。
  ⑦ business 一字不动: PPT 封面深底满铺 + stats 直角 card_fill 卡; Word 深底题块 + 标题底线;
     Excel 表头主色满底。

Run with UTF-8:  python -X utf8 tests/smoke_v171.py
"""

import os
import re
import sys
import tempfile
import zipfile

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_SRC = os.path.join(_ROOT, "src")
sys.path.insert(0, _SRC)

_DATA = os.path.join(tempfile.gettempdir(), "acc_smoke_v171_data")
os.makedirs(_DATA, exist_ok=True)
os.environ["WCW_DATA_DIR"] = _DATA

import ai_computer_control.server as server  # noqa: E402
from ai_computer_control.tools import office_style as tokens  # noqa: E402

_FNS = {t.name: t.fn for t in server.mcp._tool_manager.list_tools()}
_FAILURES: list[str] = []


def check(cond: bool, msg: str):
    print(f"  [{'ok  ' if cond else 'FAIL'}] {msg}")
    if not cond:
        _FAILURES.append(msg)


def _fill_rgb(sh):
    try:
        c = sh.fill.fore_color.rgb
        return (c[0], c[1], c[2])
    except Exception:  # noqa: BLE001
        return None


def _slide_xmls(pptx_path):
    """Return the slide XML strings in slide order."""
    with zipfile.ZipFile(pptx_path) as z:
        names = [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)]
        names.sort(key=lambda n: int(re.search(r"(\d+)", n).group()))
        return [z.read(n).decode("utf-8") for n in names]


def main() -> int:
    from openpyxl import load_workbook
    from pptx import Presentation
    from pptx.util import Emu

    body_font = tokens.get_style("business")["body_font"]  # 微软雅黑 (all three styles)

    # ============================================================ ① excel_chart 轴标题
    print("== ① excel_chart 轴标题: 自动推导 / 显式 / 关闭 / 饼图忽略 / delete=False / 字体 ==")
    xp = os.path.join(_DATA, "轴标题_单系列.xlsx")
    _FNS["write_excel"](path=xp, headers=["季度", "销售额(万元)"],
                        data=[["Q1", "1250"], ["Q2", "1580"]], style="business")
    r1 = _FNS["excel_chart"](path=xp, sheet="Sheet1", chart_type="bar",
                             data_range="A1:B3", title="销售额")
    check(r1.get("success") is True, f"单系列 bar ok (got {r1.get('error')})")
    check(r1.get("x_title") == "季度", f"x_title 自动推导=首列表头「季度」(got {r1.get('x_title')!r})")
    check(r1.get("y_title") == "销售额(万元)",
          f"y_title 自动推导=单系列表头「销售额(万元)」(got {r1.get('y_title')!r})")

    xp2 = os.path.join(_DATA, "轴标题_多系列.xlsx")
    _FNS["write_excel"](path=xp2, headers=["季度", "销售额", "利润"],
                        data=[["Q1", "1250", "320"], ["Q2", "1580", "410"]], style="business")
    r2 = _FNS["excel_chart"](path=xp2, sheet="Sheet1", chart_type="line",
                             data_range="A1:C3", title="多系列")
    check(r2.get("x_title") == "季度" and r2.get("y_title") == "",
          f"多系列: x=首列表头, y=空(图例承担) (got x={r2.get('x_title')!r} y={r2.get('y_title')!r})")

    r3 = _FNS["excel_chart"](path=xp2, sheet="Sheet1", chart_type="bar", data_range="A1:C3",
                             title="显式", target_cell="H20", x_title="时间", y_title="金额(元)")
    check(r3.get("x_title") == "时间" and r3.get("y_title") == "金额(元)",
          f"显式 x_title/y_title 覆盖自动推导 (got {r3.get('x_title')!r}/{r3.get('y_title')!r})")

    r4 = _FNS["excel_chart"](path=xp2, sheet="Sheet1", chart_type="bar", data_range="A1:C3",
                             title="关闭", target_cell="H40", x_title="", y_title="")
    check(r4.get("x_title") == "" and r4.get("y_title") == "", "传 '' 显式关闭两轴标题")

    r5 = _FNS["excel_chart"](path=xp, sheet="Sheet1", chart_type="pie", data_range="A1:B3",
                             title="饼", target_cell="H20", x_title="X", y_title="Y")
    check(r5.get("success") is True and r5.get("x_title") == "" and r5.get("y_title") == "",
          "饼图忽略 x_title/y_title 且照常成功")

    # reload and assert the persisted structure on the single-series bar chart
    wb = load_workbook(xp)
    bar = [c for c in wb.active._charts if c.tagname == "barChart"][0]
    check(bar.x_axis.delete is False and bar.y_axis.delete is False,
          "delTitle 陷阱规避: x_axis.delete == False 且 y_axis.delete == False")
    xt_run = bar.x_axis.title.tx.rich.p[0].r[0]
    yt_run = bar.y_axis.title.tx.rich.p[0].r[0]
    check(xt_run.t == "季度" and yt_run.t == "销售额(万元)",
          f"轴标题文本落盘 (got x={xt_run.t!r} y={yt_run.t!r})")
    check(xt_run.rPr.latin.typeface == body_font and xt_run.rPr.ea.typeface == body_font,
          f"轴标题字体 latin+ea == {body_font}")

    # 向后兼容: 老签名 (无 x_title/y_title kwargs) — 上面 r1/r2 即老签名调用, 已验证成功。
    check(r1.get("success") is True and r2.get("success") is True,
          "向后兼容: 老签名 excel_chart 调用照常成功 (轴标题自动补上)")

    # ============================================================ ② chart_image 轴标签
    print("\n== ② chart_image 轴标签: 顶层参数 / data 内键 / 自动推导 / '' 关闭 / 饼图忽略 ==")
    c1 = _FNS["chart_image"](path=os.path.join(_DATA, "c1.png"), chart_type="bar",
                             data={"labels": ["Q1", "Q2"],
                                   "series": [{"name": "销售额", "values": [1, 2]}]}, title="单")
    check(c1.get("success") is True and c1.get("y_title") == "销售额" and c1.get("x_title") == "",
          f"单系列自动: y=系列 name, x=空 (got x={c1.get('x_title')!r} y={c1.get('y_title')!r})")
    c2 = _FNS["chart_image"](path=os.path.join(_DATA, "c2.png"), chart_type="bar",
                             data={"labels": ["Q1"], "series": [{"name": "s", "values": [1]}]},
                             title="顶", x_title="季度", y_title="金额(万元)")
    check(c2.get("x_title") == "季度" and c2.get("y_title") == "金额(万元)", "顶层 x_title/y_title 生效")
    c3 = _FNS["chart_image"](path=os.path.join(_DATA, "c3.png"), chart_type="bar",
                             data={"labels": ["Q1"], "series": [{"name": "s", "values": [1]}],
                                   "x_title": "时间", "y_title": "数值"}, title="嵌")
    check(c3.get("x_title") == "时间" and c3.get("y_title") == "数值", "data 内 x_title/y_title 生效")
    c4 = _FNS["chart_image"](path=os.path.join(_DATA, "c4.png"), chart_type="bar",
                             data={"labels": ["Q1"], "series": [{"name": "s", "values": [1]}]},
                             title="闭", x_title="", y_title="")
    check(c4.get("x_title") == "" and c4.get("y_title") == "", "传 '' 显式关闭")
    c5 = _FNS["chart_image"](path=os.path.join(_DATA, "c5.png"), chart_type="pie",
                             data={"labels": ["a", "b"], "series": [{"name": "s", "values": [1, 2]}]},
                             title="饼", x_title="X", y_title="Y")
    check(c5.get("success") is True and c5.get("x_title") == "" and c5.get("y_title") == "",
          "饼图忽略轴标签且照常成功")

    # ============================================================ 共用 PPT 样例 deck
    tri_slides = [
        {"type": "title", "title": "季度业绩汇报", "subtitle": "副标题", "date": "2026-07-05"},
        {"type": "content", "title": "核心要点", "bullets": ["要点一", "要点二", "要点三"]},
        {"type": "stats", "title": "关键指标", "items": [
            {"label": "营收", "value": "¥18.9M", "note": "同比 +33%"},
            {"label": "客户", "value": "1,240"}]},
        {"type": "closing", "title": "谢谢"},
    ]
    decks = {}
    for st in ("business", "minimal", "vibrant"):
        p = os.path.join(_DATA, f"v171_{st}.pptx")
        r = _FNS["write_pptx"](path=p, slides=tri_slides, style=st)
        if not (isinstance(r, dict) and r.get("success") is True):
            check(False, f"write_pptx[{st}] ok (got {r})")
            continue
        decks[st] = p

    # ============================================================ ③ minimal PPT 版式
    print("\n== ③ minimal「墨白极简」PPT: 白底封面 + 青竖线 / 无横幅 content / 无填充 stats 卡 ==")
    mt = tokens.get_style("minimal")
    prs_m = Presentation(decks["minimal"])
    w, h = prs_m.slide_width, prs_m.slide_height
    cov = prs_m.slides[0]
    check(any(_fill_rgb(s) == (255, 255, 255) and s.width == w and s.height == h
              for s in cov.shapes), "minimal 封面底 = 白色满铺 (不再深底)")
    teal = tokens.rgb_tuple(mt["accent"])
    check(any(_fill_rgb(s) == teal and s.width < Emu(int(0.15 * 914400))
              and s.height > Emu(int(1.0 * 914400)) for s in cov.shapes),
          "minimal 封面左侧 4pt 青色竖线存在 (细而高的青色形状)")
    # content 无主色横幅: 不存在 primary(222222) 填充的全宽横条
    ink = tokens.rgb_tuple(mt["primary"])
    content_sl = prs_m.slides[1]
    check(not any(_fill_rgb(s) == ink and s.width == w for s in content_sl.shapes),
          "minimal content 无主色满铺横幅 (bigtext 版式)")
    check(any(_fill_rgb(s) == teal for s in content_sl.shapes),
          "minimal content 标题下细青线存在")
    # stats 无填充卡: 不存在 card_fill 面板
    cf = tokens.rgb_tuple(mt["card_fill"])
    stats_sl = prs_m.slides[2]
    check(not any(_fill_rgb(s) == cf for s in stats_sl.shapes),
          "minimal stats 卡片无填充 (无 card_fill 面板)")
    check(sum(1 for s in stats_sl.shapes if _fill_rgb(s) == teal) >= 2,
          "minimal stats 每卡顶部 2pt 青线 (2 张卡 ≥2 条青线)")

    # ============================================================ ④ vibrant PPT 版式
    print("\n== ④ vibrant「活力现代」PPT: 真圆角 stats 卡 (roundRect) / 封面几何装饰 / 珊瑚圆点 ==")
    xmls = _slide_xmls(decks["vibrant"])
    stats_xml = xmls[2]
    check('prst="roundRect"' in stats_xml,
          "vibrant stats 卡 prstGeom == roundRect (真圆角)")
    check(stats_xml.count('prst="roundRect"') >= 2,
          f"两张卡都是圆角 (roundRect ×{stats_xml.count('prst=\"roundRect\"')})")
    check('fmla="val 12000"' in stats_xml,
          "圆角 adjustment ≈ 0.12 (val 12000) 已显式设置")
    cover_xml = xmls[0]
    check(cover_xml.count('prst="ellipse"') >= 2 and 'prst="roundRect"' in cover_xml,
          "vibrant 封面右下几何装饰: ≥2 圆 + 1 圆角矩形")
    vt = tokens.get_style("vibrant")
    prs_v = Presentation(decks["vibrant"])
    check(any(_fill_rgb(s) == tokens.rgb_tuple(vt["title_bg"])
              and s.width == prs_v.slide_width and s.height == prs_v.slide_height
              for s in prs_v.slides[0].shapes), "vibrant 封面深靛底满铺仍在")
    check('prst="ellipse"' in xmls[1], "vibrant content 标题栏右端珊瑚圆点 (ellipse) 存在")

    # ============================================================ ⑤ Word 版式
    print("\n== ⑤ Word: minimal 白底封面+左竖线标题 / vibrant 珊瑚粗条+深靛题块 ==")
    cover = {"title": "如意 季度经营报告", "subtitle": "2026 财年第一季度",
             "date": "2026-07-05", "author": "把关人"}
    content = "# 一、概述\n正文段落。\n## 1.1 小节\n- 要点一\n"
    docs = {}
    for st in ("business", "minimal", "vibrant"):
        p = os.path.join(_DATA, f"v171_{st}.docx")
        r = _FNS["write_document"](path=p, content=content, title="年度报告", style=st,
                                   cover=cover, page_numbers=True)
        check(isinstance(r, dict) and r.get("success") is True, f"write_document[{st}] ok")
        with zipfile.ZipFile(p) as z:
            docs[st] = z.read("word/document.xml").decode("utf-8", "ignore")

    mm = docs["minimal"]
    check('w:fill="222222"' not in mm, "minimal 封面无深底题块 (白底装帧)")
    check('<w:left w:val="single" w:sz="18"' in mm and "3B7C8C" in mm,
          "minimal 正文标题左侧青色细竖线 (w:pBdr/w:left)")
    check('w:sz="12"' in mm, "minimal 封面标题下细青线 (1.5pt)")
    vv = docs["vibrant"]
    check('w:fill="F97066"' in vv, "vibrant 封面题块下沿珊瑚粗条 (shd F97066)")
    check('w:val="432"' in vv, "珊瑚粗条 0.3in exact 行高 (432 twips)")
    check('w:fill="3730A3"' in vv, "vibrant 深靛题块仍在")

    # ============================================================ ⑥ Excel 表头版式
    print("\n== ⑥ Excel: minimal 表头白底墨黑+青色下边框 (双通道) / 其它风格主色满底 ==")
    headers = ["区域", "销售额", "同比增长"]
    data = [["华东区", "1250000", "12%"], ["华南区", "980000", "-5%"]]
    for st, want_fill in (("business", "FF2F5597"), ("vibrant", "FF4F46E5")):
        p = os.path.join(_DATA, f"v171_{st}.xlsx")
        _FNS["write_excel"](path=p, headers=headers, data=data, style=st)
        _FNS["excel_beautify"](path=p, style=st)
        hc = load_workbook(p).active.cell(1, 1)
        check(hc.fill.fill_type == "solid" and hc.fill.fgColor.rgb == want_fill,
              f"{st} 表头主色满底 ({want_fill})")
    pm = os.path.join(_DATA, "v171_minimal.xlsx")
    _FNS["write_excel"](path=pm, headers=headers, data=data, style="minimal")
    raw = load_workbook(pm).active.cell(1, 1)
    check(raw.border.bottom.style == "medium" and raw.border.bottom.color.rgb == "FF3B7C8C",
          "minimal write_excel 落笔即带青色 medium 下边框")
    _FNS["excel_beautify"](path=pm, style="minimal")
    wbm = load_workbook(pm)
    hm = wbm.active.cell(1, 1)
    check(hm.fill.fill_type != "solid", "minimal beautify 后表头无满底色")
    check(hm.font.bold is True and hm.font.color.rgb == "FF222222", "minimal 表头墨黑粗体")
    check(hm.border.bottom.style == "medium" and hm.border.bottom.color.rgb == "FF3B7C8C",
          "minimal beautify 后表头底部 2pt(medium) 青色下边框")
    # idempotent: 重跑 beautify 不回退表头版式
    _FNS["excel_beautify"](path=pm, style="minimal")
    hm2 = load_workbook(pm).active.cell(1, 1)
    check(hm2.fill.fill_type != "solid" and hm2.border.bottom.style == "medium",
          "minimal beautify 重跑幂等 (表头版式不回退)")

    # ============================================================ ⑦ business 一字不动
    print("\n== ⑦ business 定稿默认一字不动: PPT 深底封面+直角卡 / Word 深底题块+底线 ==")
    bt = tokens.get_style("business")
    prs_b = Presentation(decks["business"])
    check(any(_fill_rgb(s) == tokens.rgb_tuple(bt["title_bg"])
              and s.width == prs_b.slide_width and s.height == prs_b.slide_height
              for s in prs_b.slides[0].shapes), "business PPT 封面深底满铺不变")
    b_stats_xml = _slide_xmls(decks["business"])[2]
    check('prst="roundRect"' not in b_stats_xml, "business stats 卡仍是直角 (无 roundRect)")
    check(sum(1 for s in prs_b.slides[2].shapes
              if _fill_rgb(s) == tokens.rgb_tuple(bt["card_fill"])) == 2,
          "business stats 卡 card_fill 面板不变")
    bb = docs["business"]
    check('w:fill="1F3864"' in bb, "business Word 深底题块不变")
    check('<w:bottom w:val="single" w:sz="18"' in bb, "business Word 标题底线不变")
    check('<w:left w:val="single" w:sz="18"' not in bb, "business Word 无左竖线 (未被 minimal 版式污染)")

    # ---------------------------------------------------------------- summary
    print()
    if _FAILURES:
        print(f"FAILED: {len(_FAILURES)} assertion(s)")
        for f in _FAILURES:
            print("  -", f)
        print("ACC-V171 SMOKE: FAIL")
        return 1
    print("样例(把关人 COM 导图过审美关):")
    for st in ("business", "minimal", "vibrant"):
        print(f"  {st:9s} pptx →", os.path.join(_DATA, f"v171_{st}.pptx"))
    print("ACC-V171 SMOKE: ALL PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
