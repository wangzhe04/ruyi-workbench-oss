"""PowerPoint generation (write_pptx) — v1.7 「Office 体系 2.0」(模板驱动).

Builds a 16:9 .pptx from a list of slide specs, styled entirely from the design tokens (no binary
template). python-pptx is an OPTIONAL offline dependency: the import is guarded at MODULE TOP so an
absent library never stops the whole MCP from starting (v1.4 P0 lesson) — instead the tool degrades to
a 中文人话 install hint.

v1.7: content 正文区顶上移 (1.75in→1.4in) so the vertically-centred body sits at the optical centre
(把关人 实测「重心略低」). The content title-bar underline switched from `accent` to `accent_on_light`
because business 的 accent 变成鎏金 (C9A860) 后在白底对比不足 (2.27 < 3.0); the cover/closing 金色下划线
(drawn on the dark title_bg) still use `accent` (鎏金 vs 藏蓝 = 5.12, 极佳).

中文字体: font family names ('微软雅黑') are written verbatim into the .pptx; PowerPoint resolves the
real CJK font on the (中文 Windows) target machine at open time.

Slide specs (list of dicts):
  {type:'title',   title, subtitle?, date?}        — full-bleed dark cover (v1.6.1 redesign)
  {type:'content', title, bullets:[str | {text, level}]}  — smart sizing / 2-col by count
  {type:'stats',   title, items:[{label, value, note?}]}  — number-highlight cards (v1.6.1 new)
  {type:'table',   title, headers:[...], rows:[[...], ...]}
  {type:'image',   title, image_path, caption?}
  {type:'closing', title?, subtitle?}              — full-bleed dark closing (v1.6.1 new)

v1.6.1 视觉升级: the title cover is now a full-bleed dark (title_bg) slide with a white bold title,
an accent underline and an optional bottom-right date; content slides size their text and auto-split
into two columns by bullet count and vertically centre the body; a new `stats` layout renders even,
number-first cards; a new `closing` layout mirrors the cover. All colours come from office_style
tokens (title_bg / on_dark_subtle / card_fill new in v1.6.1) — no hard-coded hex in this module.

v1.7.1 版式个性 (用户反馈「方案二/三只是换了配色的方案一」): layouts now branch on the office_style
版式选择器 tokens (pptx_cover / pptx_content_title / pptx_stats), giving each style its own design
language while business (定稿默认) keeps the v1.6.1 layout byte-for-byte:
  * minimal「墨白极简」— 高级文印/咨询装帧: WHITE cover (no more dark full-bleed) with a huge
    left-aligned ink-black title and a 4pt teal vertical line through the title band; content titles
    are big dark text over a thin teal rule (no colour bar); stats cards are unfilled — just a 2pt
    teal top rule + big ink-black number.
  * vibrant「活力现代」— 新锐发布会: deep-indigo cover with a restrained coral/light-indigo geometric
    cluster in the bottom-right (rounded square + two circles); stats cards are REAL rounded
    rectangles (roundRect, adjustments[0]=0.12, shadow.inherit=False — v1.6.1 曾留「直角更稳」注,
    这次把圆角做对); the content title bar keeps the colour band and gains a small coral end-dot.
"""

import os

from ai_computer_control.server import mcp
from ai_computer_control.tools.safety import protected_path_reason
from ai_computer_control.tools import office_style as style_tokens

# Guarded optional import (module top): absent python-pptx must NOT crash server import.
try:
    import pptx  # type: ignore  # noqa: F401
    _AVAILABLE = True
    _IMPORT_ERROR = ""
except Exception as e:  # noqa: BLE001 — optional dependency
    _AVAILABLE = False
    _IMPORT_ERROR = str(e)


def _unavailable() -> dict:
    return {"error": "PPTX 生成需要 python-pptx。离线包已含，可运行 installer 重装；或 pip install python-pptx",
            "detail": _IMPORT_ERROR}


def _protected_read_guard(path: str):
    """Image paths are READ. Refuse reading OUT of a protected tree only if it would be odd — actually
    reading is safe; but per the contract we run image_path through the protected guard on the read
    side to stay consistent with the write族 discipline. A protected image path is refused."""
    reason = protected_path_reason(path)
    if reason:
        return {"error": f"拒绝读取图片：{reason}。"}
    return None


def _protected_write_guard(path: str, allow_protected: bool):
    reason = protected_path_reason(path)
    if reason and not allow_protected:
        return {"error": f"拒绝写入：目标 {reason}。如确需写入，请传 allow_protected=true。"}
    return None


@mcp.tool(audit=True)
def write_pptx(
    path: str,
    slides: list[dict],
    style: str = "business",
    allow_protected: bool = False,
) -> dict:
    """Create a 16:9 PowerPoint (.pptx) from slide specs, styled from design tokens (模板驱动).

    slides is a list of dicts; each dict's 'type' selects the layout:
      * {'type': 'title',   'title': str, 'subtitle': str?, 'date': str?}
            — a full-bleed dark COVER:主色满铺 background, big white bold title (44–54pt by length),
              light subtitle, an accent underline, and an optional bottom-right date. No page number.
      * {'type': 'content', 'title': str, 'bullets': [str | {'text': str, 'level': int}]}
            — a bulleted content slide. Bullets may be plain strings or {text, level} for indent
              levels (0 = top). Font size and layout adapt to the bullet count: ≤3 large / 4–5 medium
              (both vertically centred), 6–10 auto-split into two columns, >10 truncated to 10 with a
              「…(内容过多，建议拆页)」 note.
      * {'type': 'stats',   'title': str, 'items': [{'label': str, 'value': str, 'note': str?}]}
            — number-highlight cards (2–6 items): even cards (2–3 per row) with a big bold primary
              value, a label above and an optional grey note below. Empty / >6 items -> error.
      * {'type': 'table',   'title': str, 'headers': [str,...], 'rows': [[...], ...]}
            — a slide with a styled table (token header fill + white bold text + zebra rows).
      * {'type': 'image',   'title': str, 'image_path': str, 'caption': str?}
            — a slide showing an image (auto-fit) with an optional caption. image_path is read and
              passes the protected-path guard.
      * {'type': 'closing', 'title': str?, 'subtitle': str?}
            — a full-bleed dark CLOSING slide mirroring the cover ('谢谢' if no title). No page number.
    Content / stats / table / image slides carry a footer page number; title & closing do not.

    Args:
        path: Output .pptx path.
        slides: List of slide spec dicts (see above). Empty -> error.
        style: 'business' (default) | 'minimal' | 'vibrant'. Unknown -> 'business'.
        allow_protected: Override the protected-system-root guard on the destination (default off).

    Returns:
        dict with 'success', 'path', 'output_path', 'slides' (count), 'style'. Missing python-pptx ->
        {'error': install guidance}. Bad input -> {'error': <中文人话>}.
    """
    if not _AVAILABLE:
        return _unavailable()
    if not str(path).lower().endswith(".pptx"):
        return {"error": "path 必须以 .pptx 结尾"}
    if not isinstance(slides, list) or len(slides) == 0:
        return {"error": "slides 不能为空，至少需要一页（list of dict）"}
    guard = _protected_write_guard(path, allow_protected)
    if guard:
        return guard

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    tokens = style_tokens.get_style(style)
    resolved_style = style if style in style_tokens.STYLES else style_tokens.DEFAULT_STYLE

    primary = RGBColor(*style_tokens.rgb_tuple(tokens["primary"]))
    accent = RGBColor(*style_tokens.rgb_tuple(tokens["accent"]))          # ON DARK (cover/closing 金线)
    # v1.7: content 标题栏下划线画在白底 —— business accent 变鎏金后在白底对比不足(2.27), 改用 accent_on_light
    accent_on_light = RGBColor(*style_tokens.rgb_tuple(
        tokens.get("accent_on_light", tokens["accent"])))
    header_font_color = RGBColor(*style_tokens.rgb_tuple(tokens["header_font_color"]))
    zebra = RGBColor(*style_tokens.rgb_tuple(tokens["zebra_fill"]))
    text_color = RGBColor(*style_tokens.rgb_tuple(tokens["text_color"]))
    subtle = RGBColor(*style_tokens.rgb_tuple(tokens["subtle_color"]))
    # v1.6.1 cover / stats tokens (fall back to sensible existing tokens if a style lacks them)
    title_bg = RGBColor(*style_tokens.rgb_tuple(tokens.get("title_bg", tokens["primary"])))
    on_dark_subtle = RGBColor(*style_tokens.rgb_tuple(tokens.get("on_dark_subtle", tokens["header_font_color"])))
    card_fill = RGBColor(*style_tokens.rgb_tuple(tokens.get("card_fill", tokens["zebra_fill"])))
    white = RGBColor(0xFF, 0xFF, 0xFF)
    body_font = tokens["body_font"]
    title_font = tokens["title_font"]
    # v1.7.1 版式个性选择器 + 装饰色 (缺省回落 business 版式，故老 style 不受影响)
    cover_layout = tokens.get("pptx_cover", "dark_center")
    content_title_layout = tokens.get("pptx_content_title", "bar")
    stats_layout = tokens.get("pptx_stats", "top_bar")
    accent_dot = bool(tokens.get("pptx_accent_dot", False))
    deco_1 = RGBColor(*style_tokens.rgb_tuple(tokens.get("deco_1", tokens["accent"])))
    deco_2 = RGBColor(*style_tokens.rgb_tuple(tokens.get("deco_2", tokens["primary"])))

    # 16:9 canvas.
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    def _set_run(run, size_pt, color, font_name, bold=False):
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color
        run.font.name = font_name
        run.font.bold = bold
        # East-Asian font: python-pptx only sets latin by default; set the EA typeface too so CJK
        # glyphs use our font instead of the theme's minor-EA fallback.
        try:
            rpr = run._r.get_or_add_rPr()
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            # remove any existing <a:ea> then append ours so CJK uses our typeface
            for existing in rpr.findall(ns + "ea"):
                rpr.remove(existing)
            rpr.append(rpr.makeelement(ns + "ea", {"typeface": font_name}))
        except Exception:
            pass

    def _add_footer(slide, page_no):
        box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(page_no)
        _set_run(run, 10, subtle, body_font, bold=False)

    try:
        # Validate image paths up-front (fail before writing anything).
        for i, spec in enumerate(slides):
            if not isinstance(spec, dict):
                return {"error": f"第 {i + 1} 页不是 dict：{spec!r}"}
            stype = str(spec.get("type", "")).lower()
            if stype not in ("title", "content", "stats", "table", "image", "closing"):
                return {"error": f"第 {i + 1} 页 type 非法：{spec.get('type')!r}，仅支持 "
                                 "title|content|stats|table|image|closing"}
            if stype == "stats":
                items = spec.get("items")
                if not isinstance(items, list) or len(items) == 0:
                    return {"error": f"第 {i + 1} 页 stats 缺少 items（需要 2-6 个 {{label, value, note?}}）"}
                if len(items) > 6:
                    return {"error": f"第 {i + 1} 页 stats 的 items 过多（{len(items)} 个），最多 6 个，请拆页"}
                for j, it in enumerate(items):
                    if not isinstance(it, dict):
                        return {"error": f"第 {i + 1} 页 stats 第 {j + 1} 个卡片不是 dict：{it!r}"}
            if stype == "image":
                img = spec.get("image_path")
                if not img or not os.path.exists(img):
                    return {"error": f"第 {i + 1} 页图片不存在：{img!r}"}
                g = _protected_read_guard(img)
                if g:
                    return g

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]  # fully blank layout — we place everything ourselves.

        for page_no, spec in enumerate(slides, 1):
            stype = str(spec.get("type", "")).lower()
            slide = prs.slides.add_slide(blank)

            if stype == "title":
                if cover_layout == "light_left":
                    _cover_light_left(slide, spec, SLIDE_W, SLIDE_H,
                                      text_color=text_color, accent=accent, subtle=subtle,
                                      title_font=title_font, body_font=body_font, set_run=_set_run,
                                      is_closing=False)
                elif cover_layout == "dark_deco":
                    _cover_dark_deco(slide, spec, SLIDE_W, SLIDE_H,
                                     title_bg=title_bg, white=white, accent=accent,
                                     on_dark_subtle=on_dark_subtle, deco_1=deco_1, deco_2=deco_2,
                                     title_font=title_font, body_font=body_font, set_run=_set_run,
                                     is_closing=False)
                else:
                    _cover_dark_center(slide, spec, SLIDE_W, SLIDE_H,
                                       title_bg=title_bg, white=white, accent=accent,
                                       on_dark_subtle=on_dark_subtle, title_font=title_font,
                                       body_font=body_font, set_run=_set_run, is_closing=False)

            elif stype == "closing":
                if cover_layout == "light_left":
                    _cover_light_left(slide, spec, SLIDE_W, SLIDE_H,
                                      text_color=text_color, accent=accent, subtle=subtle,
                                      title_font=title_font, body_font=body_font, set_run=_set_run,
                                      is_closing=True)
                elif cover_layout == "dark_deco":
                    _cover_dark_deco(slide, spec, SLIDE_W, SLIDE_H,
                                     title_bg=title_bg, white=white, accent=accent,
                                     on_dark_subtle=on_dark_subtle, deco_1=deco_1, deco_2=deco_2,
                                     title_font=title_font, body_font=body_font, set_run=_set_run,
                                     is_closing=True)
                else:
                    _cover_dark_center(slide, spec, SLIDE_W, SLIDE_H,
                                       title_bg=title_bg, white=white, accent=accent,
                                       on_dark_subtle=on_dark_subtle, title_font=title_font,
                                       body_font=body_font, set_run=_set_run, is_closing=True)

            elif stype == "content":
                _add_title_bar(slide, spec.get("title", ""), primary, accent_on_light, title_font,
                               header_font_color, _set_run, layout=content_title_layout,
                               text_color=text_color, accent_dot_color=(accent if accent_dot else None))
                bullets = spec.get("bullets", []) or []
                if not isinstance(bullets, list):
                    return {"error": f"第 {page_no} 页 bullets 必须是列表"}

                # Normalize to (text, level) tuples up-front so layout logic is count-driven.
                norm: list[tuple[str, int]] = []
                for item in bullets:
                    if isinstance(item, dict):
                        btext = str(item.get("text", ""))
                        level = int(item.get("level", 0) or 0)
                    else:
                        btext = str(item)
                        level = 0
                    norm.append((btext, max(0, min(level, 4))))

                # >10: truncate to 10 and append a 拆页 hint as the last item.
                if len(norm) > 10:
                    norm = norm[:10]
                    norm[-1] = ("…（内容过多，建议拆页）", 0)

                n = len(norm)
                # Body area with wider left/right margins (0.8in) and clear of the title bar.
                # v1.7: 把关人 实测 content 重心略低 —— raise the body area top ~0.35in (1.75→1.4) so the
                # vertically-centred block sits a touch higher on the slide (optical centre, not geometric).
                area_left, area_right = Inches(0.8), Inches(0.8)
                area_top, area_bottom = Inches(1.4), Inches(0.9)   # top clears bar; bottom clears footer
                area_w = SLIDE_W - area_left - area_right
                area_h = SLIDE_H - area_top - area_bottom

                if n <= 3:
                    size, spacing = 24, 1.4
                    _fill_bullets(slide, norm, area_left, area_top, area_w, area_h,
                                  size, spacing, text_color, body_font, _set_run,
                                  anchor=MSO_ANCHOR.MIDDLE)
                elif n <= 5:
                    size, spacing = 20, 1.3
                    _fill_bullets(slide, norm, area_left, area_top, area_w, area_h,
                                  size, spacing, text_color, body_font, _set_run,
                                  anchor=MSO_ANCHOR.MIDDLE)
                else:
                    # 6–10: two even columns, vertically centred, smaller type.
                    size, spacing = 17, 1.25
                    half = (n + 1) // 2
                    left_items, right_items = norm[:half], norm[half:]
                    gutter = Inches(0.5)
                    col_w = (area_w - gutter) // 2
                    _fill_bullets(slide, left_items, area_left, area_top, col_w, area_h,
                                  size, spacing, text_color, body_font, _set_run,
                                  anchor=MSO_ANCHOR.MIDDLE)
                    _fill_bullets(slide, right_items, area_left + col_w + gutter, area_top, col_w,
                                  area_h, size, spacing, text_color, body_font, _set_run,
                                  anchor=MSO_ANCHOR.MIDDLE)

            elif stype == "stats":
                _add_title_bar(slide, spec.get("title", ""), primary, accent_on_light, title_font,
                               header_font_color, _set_run, layout=content_title_layout,
                               text_color=text_color, accent_dot_color=(accent if accent_dot else None))
                items = spec.get("items") or []   # already validated (2-6 dicts) up-front
                _render_stats_cards(
                    slide, items, SLIDE_W, SLIDE_H,
                    primary=primary, card_fill=card_fill, subtle=subtle, body_font=body_font,
                    title_font=title_font, text_color=text_color, accent=accent,
                    layout=stats_layout, set_run=_set_run, Inches=Inches, Pt=Pt,
                    PP_ALIGN=PP_ALIGN, MSO_ANCHOR=MSO_ANCHOR)

            elif stype == "table":
                _add_title_bar(slide, spec.get("title", ""), primary, accent_on_light, title_font,
                               header_font_color, _set_run, layout=content_title_layout,
                               text_color=text_color, accent_dot_color=(accent if accent_dot else None))
                headers = spec.get("headers", []) or []
                rows = spec.get("rows", []) or []
                if not headers:
                    return {"error": f"第 {page_no} 页 table 缺少 headers"}
                n_cols = len(headers)
                n_rows = len(rows) + 1  # + header
                tbl_shape = slide.shapes.add_table(
                    n_rows, n_cols, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.4 * n_rows))
                table = tbl_shape.table
                # header row
                for c, h in enumerate(headers):
                    cell = table.cell(0, c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = primary
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf = cell.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = str(h)
                    _set_run(run, 14, header_font_color, body_font, bold=True)
                # body rows with zebra
                for r, row in enumerate(rows, start=1):
                    striped = (r % 2 == 0)
                    for c in range(n_cols):
                        cell = table.cell(r, c)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = zebra if striped else RGBColor(0xFF, 0xFF, 0xFF)
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                        val = row[c] if c < len(row) else ""
                        tf = cell.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = str(val)
                        _set_run(run, 12, text_color, body_font, bold=False)

            elif stype == "image":
                _add_title_bar(slide, spec.get("title", ""), primary, accent_on_light, title_font,
                               header_font_color, _set_run, layout=content_title_layout,
                               text_color=text_color, accent_dot_color=(accent if accent_dot else None))
                img = spec.get("image_path")
                caption = spec.get("caption")
                # Fit within a box, preserving aspect ratio via python-pptx auto-height.
                max_w = Inches(11.5)
                left = Inches(0.9)
                top = Inches(1.8)
                try:
                    pic = slide.shapes.add_picture(img, left, top, width=max_w)
                    # If it overflows vertically, re-add constrained by height instead.
                    if pic.height > Inches(4.8):
                        # remove and re-add by height
                        sp = pic._element
                        sp.getparent().remove(sp)
                        pic = slide.shapes.add_picture(img, left, top, height=Inches(4.8))
                except Exception as e:  # noqa: BLE001
                    return {"error": f"第 {page_no} 页插入图片失败：{e}"}
                if caption:
                    cbox = slide.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4))
                    ctf = cbox.text_frame
                    ctf.word_wrap = True
                    cp = ctf.paragraphs[0]
                    cp.alignment = PP_ALIGN.CENTER
                    crun = cp.add_run()
                    crun.text = str(caption)
                    _set_run(crun, 12, subtle, body_font, bold=False)

            # Cover & closing are full-bleed dark and carry no page number.
            if stype not in ("title", "closing"):
                _add_footer(slide, page_no)

        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        prs.save(path)
        return {
            "success": True,
            "path": os.path.abspath(path),
            "output_path": os.path.abspath(path),
            "slides": len(slides),
            "style": resolved_style,
        }
    except Exception as e:  # noqa: BLE001
        return {"error": f"PPTX 生成失败：{e}"}


def _add_title_bar(slide, title_text, primary, accent, title_font, header_font_color, set_run,
                   *, layout="bar", text_color=None, accent_dot_color=None):
    """Shared slide-title header. Two layouts:

      * 'bar' (business / vibrant) — a full-width primary bar with white bold title + a thin accent
        underline. vibrant additionally drops a small珊瑚 accent dot at the bar's right end (via
        accent_dot_color) as a发布会-style finishing flourish.
      * 'bigtext' (minimal) — NO colour block; a large near-black title left-aligned on the白底,
        with a thin accent teal hairline just below it (editorial / 咨询装帧 feel).
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    if layout == "bigtext":
        # minimal: no bar. Big dark title on white + thin teal rule under it.
        tbox = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.8), Inches(0.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(title_text)
        set_run(run, 30, text_color if text_color is not None else header_font_color,
                title_font, bold=True)
        # thin teal hairline under the title (short, left-anchored — a rule, not a full bar)
        rule = slide.shapes.add_shape(1, Inches(0.8), Inches(1.32), Inches(2.4), Pt(2.5))
        rule.fill.solid()
        rule.fill.fore_color.rgb = accent
        rule.line.fill.background()
        rule.shadow.inherit = False
        return

    # 'bar' (default): primary bar with white bold title + thin accent underline.
    bar = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(1.25))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(title_text)
    set_run(run, 26, header_font_color, title_font, bold=True)
    # thin accent underline just below the bar
    line = slide.shapes.add_shape(1, 0, Inches(1.25), Inches(13.333), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # vibrant: 珊瑚 accent dot at the right end of the bar (发布会 finishing flourish).
    if accent_dot_color is not None:
        d = Inches(0.34)
        dot = slide.shapes.add_shape(9, Inches(13.333) - Inches(0.7), int((Inches(1.25) - d) / 2),
                                     d, d)  # 9 = OVAL
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_dot_color
        dot.line.fill.background()
        dot.shadow.inherit = False


def _cover_dark_center(slide, spec, slide_w, slide_h, *, title_bg, white, accent, on_dark_subtle,
                       title_font, body_font, set_run, is_closing):
    """business cover/closing (v1.6.1 原版式，不动): full-bleed dark, big white centred title, a金色
    accent underline, light subtitle, optional bottom-right date. No page number."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    _fill_slide_bg(slide, title_bg, slide_w, slide_h)

    if is_closing:
        title_text = str(spec.get("title") or "谢谢")
        t_top, tsize, u_top, sub_top = Inches(2.7), 48, Inches(4.45), Inches(4.7)
    else:
        title_text = str(spec.get("title", ""))
        tsize = 54 if len(title_text) <= 12 else (48 if len(title_text) <= 22 else 44)
        t_top, u_top, sub_top = Inches(2.5), Inches(4.35), Inches(4.6)

    tbox = slide.shapes.add_textbox(Inches(1.0), t_top, Inches(11.333), Inches(1.7))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    set_run(run, tsize, white, title_font, bold=True)

    uw = Inches(1.5)
    uline = slide.shapes.add_shape(1, int((slide_w - uw) / 2), u_top, uw, Pt(5))
    uline.fill.solid()
    uline.fill.fore_color.rgb = accent
    uline.line.fill.background()

    subtitle = spec.get("subtitle")
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.0), sub_top, Inches(11.333), Inches(0.9))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = str(subtitle)
        set_run(srun, 20, on_dark_subtle, body_font, bold=False)

    if not is_closing:
        date_text = spec.get("date")
        if date_text:
            dbox = slide.shapes.add_textbox(Inches(7.0), Inches(6.75), Inches(5.6), Inches(0.4))
            dtf = dbox.text_frame
            dtf.word_wrap = False
            dp = dtf.paragraphs[0]
            dp.alignment = PP_ALIGN.RIGHT
            drun = dp.add_run()
            drun.text = str(date_text)
            set_run(drun, 12, on_dark_subtle, body_font, bold=False)


def _cover_light_left(slide, spec, slide_w, slide_h, *, text_color, accent, subtle,
                      title_font, body_font, set_run, is_closing):
    """minimal 「墨白极简」cover/closing — 高级文印/咨询装帧: 纯白底 (反转，不再深底满铺), 超大墨黑
    左对齐标题 (56–60pt), 左侧一条 4pt 青色竖线贯穿标题区, 右下极小页脚字。closing 同构 ('谢谢' 默认)。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # white full-bleed background (explicit — so it reads as intentional white, not theme default).
    _fill_slide_bg(slide, RGBColor_white(), slide_w, slide_h)

    if is_closing:
        title_text = str(spec.get("title") or "谢谢")
    else:
        title_text = str(spec.get("title", ""))

    # left teal vertical line spanning the title band (4pt wide).
    line_top, line_h = Inches(2.3), Inches(2.4)
    vline = slide.shapes.add_shape(1, Inches(0.9), line_top, Pt(4), line_h)
    vline.fill.solid()
    vline.fill.fore_color.rgb = accent
    vline.line.fill.background()
    vline.shadow.inherit = False

    # super-large near-black left-aligned title, just right of the line.
    tsize = 60 if len(title_text) <= 10 else (56 if len(title_text) <= 18 else 44)
    tbox = slide.shapes.add_textbox(Inches(1.25), line_top, Inches(11.2), line_h)
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_run(run, tsize, text_color, title_font, bold=True)

    subtitle = spec.get("subtitle")
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.28), line_top + line_h + Inches(0.1),
                                        Inches(11.0), Inches(0.7))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        srun = sp.add_run()
        srun.text = str(subtitle)
        set_run(srun, 18, subtle, body_font, bold=False)

    # tiny bottom-right footer text (date on cover; subtle brand line otherwise).
    footer_text = None
    if not is_closing:
        footer_text = spec.get("date")
    if footer_text:
        fbox = slide.shapes.add_textbox(Inches(7.0), Inches(6.95), Inches(5.6), Inches(0.35))
        ftf = fbox.text_frame
        ftf.word_wrap = False
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        frun = fp.add_run()
        frun.text = str(footer_text)
        set_run(frun, 10, subtle, body_font, bold=False)


def _cover_dark_deco(slide, spec, slide_w, slide_h, *, title_bg, white, accent, on_dark_subtle,
                     deco_1, deco_2, title_font, body_font, set_run, is_closing):
    """vibrant 「活力现代」cover/closing — 新锐发布会: 深靛底 + 右下角一组珊瑚/浅靛几何装饰 (两三个错落
    的实心圆/圆角矩形, 勿花哨过度) + 居中白粗标题 + 金/珊瑚短下划线 + 浅靛副题。closing 同构。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    _fill_slide_bg(slide, title_bg, slide_w, slide_h)

    # --- bottom-right geometric cluster (drawn before text so text/underline sit above it) ---
    # a big lighter-indigo rounded square, an overlapping coral circle, and a small coral circle —
    # 错落 (offset) but restrained. All shadow-less, no outline.
    big = slide.shapes.add_shape(5, slide_w - Inches(2.7), slide_h - Inches(2.7),
                                 Inches(2.2), Inches(2.2))  # 5 = ROUNDED_RECTANGLE
    _round_adj(big, 0.18)
    big.fill.solid()
    big.fill.fore_color.rgb = deco_2
    big.line.fill.background()
    big.shadow.inherit = False

    circ = slide.shapes.add_shape(9, slide_w - Inches(1.55), slide_h - Inches(3.05),
                                  Inches(1.4), Inches(1.4))  # 9 = OVAL
    circ.fill.solid()
    circ.fill.fore_color.rgb = deco_1
    circ.line.fill.background()
    circ.shadow.inherit = False

    small = slide.shapes.add_shape(9, slide_w - Inches(3.15), slide_h - Inches(1.35),
                                   Inches(0.6), Inches(0.6))
    small.fill.solid()
    small.fill.fore_color.rgb = deco_1
    small.line.fill.background()
    small.shadow.inherit = False

    # --- title (centred white bold) ---
    if is_closing:
        title_text = str(spec.get("title") or "谢谢")
        t_top, tsize, u_top, sub_top = Inches(2.6), 48, Inches(4.35), Inches(4.6)
    else:
        title_text = str(spec.get("title", ""))
        tsize = 54 if len(title_text) <= 12 else (48 if len(title_text) <= 22 else 44)
        t_top, u_top, sub_top = Inches(2.4), Inches(4.25), Inches(4.5)

    tbox = slide.shapes.add_textbox(Inches(1.0), t_top, Inches(11.333), Inches(1.7))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    set_run(run, tsize, white, title_font, bold=True)

    uw = Inches(1.6)
    uline = slide.shapes.add_shape(1, int((slide_w - uw) / 2), u_top, uw, Pt(5))
    uline.fill.solid()
    uline.fill.fore_color.rgb = accent      # 珊瑚 accent — pops on the deep indigo
    uline.line.fill.background()
    uline.shadow.inherit = False

    subtitle = spec.get("subtitle")
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.0), sub_top, Inches(11.333), Inches(0.9))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = str(subtitle)
        set_run(srun, 20, on_dark_subtle, body_font, bold=False)

    if not is_closing:
        date_text = spec.get("date")
        if date_text:
            dbox = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(5.6), Inches(0.4))
            dtf = dbox.text_frame
            dtf.word_wrap = False
            dp = dtf.paragraphs[0]
            dp.alignment = PP_ALIGN.LEFT   # bottom-LEFT so it clears the bottom-right deco cluster
            drun = dp.add_run()
            drun.text = str(date_text)
            set_run(drun, 12, on_dark_subtle, body_font, bold=False)


def RGBColor_white():
    from pptx.dml.color import RGBColor
    return RGBColor(0xFF, 0xFF, 0xFF)


def _round_adj(shape, value):
    """Set a ROUNDED_RECTANGLE's corner-radius adjustment (adj[0]) to `value` (0..0.5). python-pptx
    exposes shape.adjustments; guard defensively as some builds/shape-types lack it."""
    try:
        shape.adjustments[0] = float(value)
    except Exception:  # noqa: BLE001
        pass


def _fill_slide_bg(slide, color, slide_w, slide_h):
    """Paint a full-bleed rectangle of `color` behind everything (cover / closing 满铺底色)."""
    bg = slide.shapes.add_shape(1, 0, 0, slide_w, slide_h)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Keep the paint at the very back so later text/underline shapes draw on top of it.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)  # index 2 == first shape slot after nvGrpSpPr / grpSpPr


def _fill_bullets(slide, items, left, top, width, height, size_pt, line_spacing,
                  text_color, body_font, set_run, anchor):
    """Fill one textbox with normalized (text, level) bullets at a fixed size / line spacing.

    items: list of (text, level). Vertically anchored via `anchor` so short lists sit centred in
    the body area rather than clinging to the top. Line spacing is a float multiplier (1.4 == 140%).
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for btext, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.add_run()
        bullet_char = "• " if level == 0 else ("– " if level == 1 else "· ")
        run.text = bullet_char + btext
        # Nested levels step down a touch so hierarchy reads without going tiny.
        set_run(run, max(12, size_pt - level * 2), text_color, body_font, bold=False)
    return box


def _render_stats_cards(slide, items, slide_w, slide_h, *, primary, card_fill, subtle,
                        body_font, title_font, text_color, accent, layout, set_run,
                        Inches, Pt, PP_ALIGN, MSO_ANCHOR):
    """Render 2–6 even number-highlight cards, vertically centred in the body area. Three looks:

      * 'top_bar' (business, 原版式不动) — light `card_fill` panel + a `primary` bar across its top,
        big bold `primary` value.
      * 'top_rule' (minimal) — NO fill; only a 2pt `accent` teal rule across the card top, big bold
        near-black (`text_color`) value. Editorial / 咨询装帧 restraint.
      * 'rounded' (vibrant) — a真圆角 (ROUNDED_RECTANGLE, adj≈0.12) `card_fill` panel with a `primary`
        top bar, big bold `primary` value. 发布会 softness (shadow explicitly off).

    Layout math (row/col placement) is identical across looks: 2–3 items -> one row; 4 -> 2×2; 5–6 ->
    two rows of 3 (last row may hold 2).
    """
    n = len(items)
    per_row = n if n <= 3 else 3
    rows = (n + per_row - 1) // per_row

    area_left, area_right = Inches(0.9), Inches(0.9)
    area_top, area_bottom = Inches(1.7), Inches(0.9)
    area_w = slide_w - area_left - area_right
    area_h = slide_h - area_top - area_bottom

    h_gap, v_gap = Inches(0.4), Inches(0.4)
    card_h = (area_h - v_gap * (rows - 1)) // rows if rows else area_h
    # Cap card height so a single row does not stretch a card into a full-slide block.
    max_card_h = Inches(2.4)
    if card_h > max_card_h:
        card_h = max_card_h
    # Vertically centre the whole card block within the body area.
    block_h = card_h * rows + v_gap * (rows - 1)
    block_top = area_top + (area_h - block_h) // 2

    bar_h = Pt(3)
    rule_h = Pt(2)
    value_color = text_color if layout == "top_rule" else primary
    for idx, it in enumerate(items):
        r = idx // per_row
        c = idx % per_row
        # cards in the last (possibly shorter) row are centred by using that row's own count.
        row_count = per_row if (r < rows - 1) else (n - per_row * (rows - 1))
        card_w = (area_w - h_gap * (row_count - 1)) // row_count
        row_block_w = card_w * row_count + h_gap * (row_count - 1)
        row_left = area_left + (area_w - row_block_w) // 2
        cx = row_left + c * (card_w + h_gap)
        cy = block_top + r * (card_h + v_gap)

        if layout == "top_rule":
            # minimal: no filled panel — just a top teal hairline (a "card" defined by its rule).
            rule = slide.shapes.add_shape(1, int(cx), int(cy), int(card_w), rule_h)
            rule.fill.solid()
            rule.fill.fore_color.rgb = accent
            rule.line.fill.background()
            rule.shadow.inherit = False
        else:
            # card panel — rounded for vibrant, square for business.
            shape_id = 5 if layout == "rounded" else 1  # 5 = ROUNDED_RECTANGLE
            panel = slide.shapes.add_shape(shape_id, int(cx), int(cy), int(card_w), int(card_h))
            if layout == "rounded":
                _round_adj(panel, 0.12)   # v1.6.1 曾留「直角更稳」注 —— 这次把圆角做对
            panel.fill.solid()
            panel.fill.fore_color.rgb = card_fill
            panel.line.fill.background()
            panel.shadow.inherit = False
            # top accent bar (primary)
            top_bar = slide.shapes.add_shape(1, int(cx), int(cy), int(card_w), bar_h)
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = primary
            top_bar.line.fill.background()
            top_bar.shadow.inherit = False

        label = str(it.get("label", ""))
        value = str(it.get("value", ""))
        note = it.get("note")

        # one textbox holds label / value / note, vertically centred inside the card.
        tbox = slide.shapes.add_textbox(int(cx), int(cy), int(card_w), int(card_h))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)

        p_label = tf.paragraphs[0]
        p_label.alignment = PP_ALIGN.CENTER
        r_label = p_label.add_run()
        r_label.text = label
        set_run(r_label, 13, subtle, body_font, bold=False)

        p_value = tf.add_paragraph()
        p_value.alignment = PP_ALIGN.CENTER
        p_value.space_before = Pt(2)
        r_value = p_value.add_run()
        r_value.text = value
        # value size steps down a little when many cards share a row.
        vsize = 40 if n <= 3 else 36
        set_run(r_value, vsize, value_color, title_font, bold=True)

        if note:
            p_note = tf.add_paragraph()
            p_note.alignment = PP_ALIGN.CENTER
            p_note.space_before = Pt(2)
            r_note = p_note.add_run()
            r_note.text = str(note)
            set_run(r_note, 11, subtle, body_font, bold=False)
